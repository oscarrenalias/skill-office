"""Shared fixtures for the pypptx test suite.

A real minimal .pptx is generated via python-pptx rather than using mocks.
A real minimal .xlsx is generated via openpyxl rather than using mocks.
"""
import datetime
import pytest
from pathlib import Path
from pptx import Presentation


@pytest.fixture
def minimal_pptx(tmp_path):
    """Create a minimal .pptx with 3 slides.

    Slide 1: Title Slide layout — has title + subtitle text.
    Slide 2: Title and Content layout — has title text.
    Slide 3: Blank layout — marked hidden via show='0'.
    """
    prs = Presentation()

    # Slide 1 — title slide with text
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Title Slide"
    slide.placeholders[1].text = "Subtitle text"

    # Slide 2 — title + content with text
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Slide Two"

    # Slide 3 — blank, hidden
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    slide._element.set("show", "0")

    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def minimal_xlsx(tmp_path):
    """Create a minimal .xlsx workbook with two sheets via openpyxl.

    Sheet 1 ("Data"): header row + 3 data rows with mixed types
      (int, float, str, bool, datetime).
    Sheet 2 ("Summary"): a single label cell.

    Returns the Path to the file inside tmp_path.
    """
    import openpyxl

    wb = openpyxl.Workbook()

    # Sheet 1 — Data
    ws1 = wb.active
    ws1.title = "Data"
    ws1.append(["Name", "Count", "Score", "Active", "Timestamp"])
    ws1.append(["Alpha", 1, 1.1, True, datetime.datetime(2025, 1, 6, 9, 0, 0)])
    ws1.append(["Beta", 42, 3.14, False, datetime.datetime(2025, 3, 1, 12, 30, 0)])
    ws1.append(["Gamma", 0, 0.0, True, datetime.datetime(2025, 6, 15, 17, 45, 0)])

    # Sheet 2 — Summary
    ws2 = wb.create_sheet("Summary")
    ws2["A1"] = "Summary sheet"

    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def unpacked_pptx(tmp_path, minimal_pptx):
    """Unpack minimal_pptx into a temporary directory and return the path."""
    from pypptx.ops.pack import unpack

    out_dir = tmp_path / "unpacked"
    unpack(minimal_pptx, out_dir)
    return out_dir
