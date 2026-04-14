"""Tests for pypptx.ops.clean — orphan detection and removal."""
import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from pypptx.ops.clean import clean_unused_files
from pypptx.ops.pack import unpack


class TestCleanUnusedFiles:
    def test_empty_result_when_nothing_orphaned(self, unpacked_pptx):
        removed = clean_unused_files(unpacked_pptx)
        assert removed == []

    def test_removes_unreferenced_media_file(self, unpacked_pptx):
        """A media file not in any .rels chain is detected and removed."""
        media_dir = unpacked_pptx / "ppt" / "media"
        media_dir.mkdir(parents=True, exist_ok=True)
        orphan = media_dir / "orphan.png"
        orphan.write_bytes(b"\x89PNG\r\n\x1a\n")
        removed = clean_unused_files(unpacked_pptx)
        assert "ppt/media/orphan.png" in removed
        assert not orphan.exists()

    def test_returns_sorted_list(self, unpacked_pptx):
        """Returned list is sorted."""
        media_dir = unpacked_pptx / "ppt" / "media"
        media_dir.mkdir(parents=True, exist_ok=True)
        (media_dir / "z_orphan.png").write_bytes(b"\x89PNG")
        (media_dir / "a_orphan.png").write_bytes(b"\x89PNG")
        removed = clean_unused_files(unpacked_pptx)
        orphans = [r for r in removed if "orphan" in r]
        assert orphans == sorted(orphans)

    def test_directory_path_input(self, unpacked_pptx):
        result = clean_unused_files(unpacked_pptx)
        assert isinstance(result, list)

    def test_pptx_file_input_round_trip(self, minimal_pptx, tmp_path):
        """Works on a .pptx path: unpacks, cleans, repacks; file is still valid."""
        pptx_copy = tmp_path / "copy.pptx"
        shutil.copy2(minimal_pptx, pptx_copy)
        removed = clean_unused_files(pptx_copy)
        assert isinstance(removed, list)
        prs = Presentation(str(pptx_copy))
        assert len(prs.slides) == 3

    def test_slide_absent_from_sldIdLst_treated_as_orphan(self, unpacked_pptx):
        """A slide file dropped directly into ppt/slides/ without updating
        presentation.xml is not reachable and must be cleaned."""
        orphan_slide = unpacked_pptx / "ppt" / "slides" / "slideOrphan.xml"
        orphan_slide.write_text(
            "<?xml version='1.0' encoding='UTF-8'?>"
            "<p:sld xmlns:p='http://schemas.openxmlformats.org/presentationml/2006/main'/>",
            encoding="UTF-8",
        )
        removed = clean_unused_files(unpacked_pptx)
        assert "ppt/slides/slideOrphan.xml" in removed
        assert not orphan_slide.exists()

    def test_reachable_core_files_retained(self, unpacked_pptx):
        """After cleaning, core package files are still present."""
        clean_unused_files(unpacked_pptx)
        after = {
            f.relative_to(unpacked_pptx).as_posix()
            for f in unpacked_pptx.rglob("*")
            if f.is_file()
        }
        assert "[Content_Types].xml" in after
        assert "ppt/presentation.xml" in after

    def test_pptx_file_input_adds_orphan_then_cleans(self, minimal_pptx, tmp_path):
        """Orphan added to a .pptx is removed and file remains openable."""
        pptx_copy = tmp_path / "copy.pptx"
        shutil.copy2(minimal_pptx, pptx_copy)

        # Inject an orphan file into the archive
        import zipfile

        orphan_name = "ppt/media/injected_orphan.png"
        with zipfile.ZipFile(pptx_copy, "a") as zf:
            zf.writestr(orphan_name, b"\x89PNG")

        removed = clean_unused_files(pptx_copy)
        assert orphan_name in removed

        prs = Presentation(str(pptx_copy))
        assert len(prs.slides) == 3
