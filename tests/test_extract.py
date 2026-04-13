"""Tests for pypptx.ops.extract — extract_text."""
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from pypptx.ops.extract import extract_text


class TestExtractText:
    def test_all_slides_included(self, minimal_pptx):
        result = extract_text(minimal_pptx)
        assert "--- Slide 1 ---" in result
        assert "--- Slide 2 ---" in result
        assert "--- Slide 3 ---" in result

    def test_slide_delimiters_in_order(self, minimal_pptx):
        result = extract_text(minimal_pptx)
        idx1 = result.index("--- Slide 1 ---")
        idx2 = result.index("--- Slide 2 ---")
        idx3 = result.index("--- Slide 3 ---")
        assert idx1 < idx2 < idx3

    def test_text_content_extracted(self, minimal_pptx):
        result = extract_text(minimal_pptx)
        assert "Title Slide" in result
        assert "Subtitle text" in result

    def test_filtered_slides_includes_only_requested(self, minimal_pptx):
        result = extract_text(minimal_pptx, slides=[1])
        assert "--- Slide 1 ---" in result
        assert "--- Slide 2 ---" not in result
        assert "--- Slide 3 ---" not in result

    def test_filtered_slides_multiple(self, minimal_pptx):
        result = extract_text(minimal_pptx, slides=[1, 3])
        assert "--- Slide 1 ---" in result
        assert "--- Slide 2 ---" not in result
        assert "--- Slide 3 ---" in result

    def test_none_slides_returns_all(self, minimal_pptx):
        result_all = extract_text(minimal_pptx, slides=None)
        result_default = extract_text(minimal_pptx)
        assert result_all == result_default

    def test_non_text_shapes_skipped(self, tmp_path):
        """Shapes without a text frame do not cause errors and are skipped."""
        prs = Presentation()
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)
        # Add a line shape (no text frame)
        slide.shapes.add_connector(
            1,  # MSO_CONNECTOR_TYPE.STRAIGHT
            Inches(0), Inches(0), Inches(1), Inches(1),
        )
        path = tmp_path / "noshapes.pptx"
        prs.save(str(path))
        result = extract_text(path)
        assert "--- Slide 1 ---" in result

    def test_empty_presentation(self, tmp_path):
        """Presentation with no slides returns an empty string."""
        prs = Presentation()
        path = tmp_path / "empty.pptx"
        prs.save(str(path))
        result = extract_text(path)
        assert result == ""

    def test_reading_order_sort(self, tmp_path):
        """Shapes are sorted by top position; upper shapes appear first."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)

        # Lower shape (large top offset) — text: "BOTTOM"
        bottom_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(3), Inches(1)
        )
        bottom_box.text_frame.text = "BOTTOM"

        # Upper shape (small top offset) — text: "TOP"
        top_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(3), Inches(1)
        )
        top_box.text_frame.text = "TOP"

        path = tmp_path / "order.pptx"
        prs.save(str(path))

        result = extract_text(path)
        top_pos = result.index("TOP")
        bottom_pos = result.index("BOTTOM")
        assert top_pos < bottom_pos
