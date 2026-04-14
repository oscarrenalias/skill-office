"""Verification checks for .pptx files."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# ── Namespaces ────────────────────────────────────────────────────────────────

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ── Check 1: unfilled placeholder hint text ───────────────────────────────────

_PLACEHOLDER_PATTERNS: tuple[str, ...] = (
    "click to add title",
    "click to add text",
    "click to add subtitle",
    "click to edit master title style",
    "click to edit master text styles",
    "add title",
    "add text",
    "enter text here",
    "place subtitle here",
)


def _check_unfilled_placeholders(
    slide_index: int,
    slide,
    errors: list[str],
) -> None:
    """Check 1: detect shapes whose text matches an unfilled placeholder pattern."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if not text.strip():
            continue
        if text.strip().lower() in _PLACEHOLDER_PATTERNS:
            errors.append(
                f"Slide {slide_index}: '{shape.name}' has unfilled placeholder"
                f' text \u2014 "{text.strip()}"'
            )


# ── Check 2: font size below minimum ─────────────────────────────────────────

_MIN_FONT_UNITS = 1200  # 12pt in hundredths-of-a-point


def _check_font_sizes(
    slide_index: int,
    slide,
    slide_height: int,
    errors: list[str],
) -> None:
    """Check 2: report fonts below 12pt; skip footer-region shapes."""
    footer_threshold = int(slide_height * 0.9)

    for shape in slide.shapes:
        # Skip shapes in the footer region (top > 90% of slide height)
        top = shape.top
        if top is not None and top > footer_threshold:
            continue

        # Walk a:rPr[@sz] elements in this shape's XML
        sp_elem = shape._element
        for rpr in sp_elem.iter(f"{{{NS_A}}}rPr"):
            sz = rpr.get("sz")
            if sz is None:
                continue
            try:
                sz_int = int(sz)
            except ValueError:
                continue
            if sz_int < _MIN_FONT_UNITS:
                # Convert from hundredths-of-a-point to pt for the message
                pt = sz_int / 100.0
                errors.append(
                    f"Slide {slide_index}: '{shape.name}' font {pt:.1f}pt"
                    f" is below minimum (12pt)"
                )


# ── Public API ────────────────────────────────────────────────────────────────


def verify_pptx(path: Path) -> dict:
    """Run quality checks on a .pptx file.

    Returns a dict with keys:
        errors      -- list of error message strings
        warnings    -- list of warning message strings
        slide_count -- total number of slides
    """
    prs = Presentation(Path(path))
    slide_height: int = prs.slide_height or Emu(6858000)  # default 7.5in

    errors: list[str] = []
    warnings: list[str] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        _check_unfilled_placeholders(slide_index, slide, errors)
        _check_font_sizes(slide_index, slide, int(slide_height), errors)

    return {
        "errors": errors,
        "warnings": warnings,
        "slide_count": len(prs.slides),
    }
